"""
PDF → PowerPoint (.pptx)
Strategy: Render each PDF page as an image, embed each image as one slide.
This is what most "PDF to PPT" tools actually do — the result preserves visual
layout perfectly but text isn't editable. For editable text, users should use
PDF to Word instead.
"""

import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches


def convert(input_path, output_path, options):
    output = output_path + ".pptx"
    dpi = int(options.get("dpi", 150))

    images = convert_from_path(input_path, dpi=dpi)

    prs = Presentation()
    # 16:9 slide size
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank

    img_dir = os.path.dirname(output_path)
    for i, img in enumerate(images):
        img_path = os.path.join(img_dir, f"slide_{i}.png")
        img.save(img_path, "PNG")

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_path, 0, 0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    prs.save(output)
    return output
